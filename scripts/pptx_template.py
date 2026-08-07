# -*- coding: utf-8 -*-
"""Generate a reusable 16:9 PPTX from the office-visual-spec design tokens.

- Titles use bold HeiTi (Microsoft YaHei / PingFang).
- Body uses SongTi serif.
- Long lists auto-paginate into multiple slides.
- Images can be inserted with add_picture().

Usage: python pptx_template.py [output.pptx]
"""
import sys
from pathlib import Path
from cli import run_main

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

PAPER = RGBColor(0xFA, 0xF8, 0xF5)
INK = RGBColor(0x22, 0x1E, 0x1A)
BODY = RGBColor(0x44, 0x3F, 0x3A)
NARRATIVE = RGBColor(0x55, 0x4F, 0x49)
MUTED = RGBColor(0x8E, 0x87, 0x7F)
RED = RGBColor(0xC6, 0x0D, 0x0D)
DEEP_RED = RGBColor(0xA5, 0x0A, 0x0A)
DIVIDER = RGBColor(0xE2, 0xDD, 0xD5)
CARD = RGBColor(0xFD, 0xF9, 0xF7)
ON_ACCENT = RGBColor(0xFA, 0xF8, 0xF5)
TITLE_FONT = "\u5fae\u8f6f\u96c5\u9ed1"
BODY_FONT = "\u5b8b\u4f53"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s_emu = min(prs.slide_width, prs.slide_height)
s_in = s_emu / 914400 / 100
s_pt = s_emu / 12700 / 100
slide_no = 0


def U(n):
    return Inches(s_in * n)


def F(n):
    return Pt(max(8, s_pt * n))


def next_page():
    global slide_no
    slide_no += 1
    return "%02d" % slide_no


def chunk(seq, size):
    return [seq[i:i + size] for i in range(0, len(seq), size)]


def set_run(run, text, size, color, bold=False, font=BODY_FONT):
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", font)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=BODY_FONT):
    box = slide.shapes.add_textbox(U(x), U(y), U(w), U(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, F(size), color, bold, font)
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, U(x), U(y), U(w), U(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def add_picture(slide, image_path, x, y, w, h):
    return slide.shapes.add_picture(str(image_path), U(x), U(y), U(w), U(h))


def add_image_slot(slide, x, y, w, h, label="\u56fe\u7247\u533a"):
    add_rect(slide, x, y, w, h, fill=None, line=DIVIDER, line_w=1.0)
    add_text(slide, x, y + h / 2 - 1.5, w, 3, label, 1.6, MUTED, align=PP_ALIGN.CENTER)


def header(slide, numeral, title, en):
    add_rect(slide, 5, 10, 5, 5, fill=INK)
    add_text(slide, 5, 10.2, 5, 4, numeral, 2.1, ON_ACCENT, align=PP_ALIGN.CENTER, font=TITLE_FONT, bold=True)
    add_text(slide, 12, 9.8, 65, 7, title, 4.2, INK, bold=True, font=TITLE_FONT)
    add_text(slide, 100, 10.5, 72, 5, en, 1.5, MUTED, align=PP_ALIGN.RIGHT, font=TITLE_FONT)
    add_rect(slide, 5, 20.5, 168, 0.25, fill=INK)


def footer(slide, page):
    add_rect(slide, 5, 91, 168, 0.13, fill=DIVIDER)
    add_text(slide, 7, 92.5, 45, 4, "\u6696\u7eb8\u793a\u4f8b\u6a21\u677f\u00b7\u8bfe\u7a0b\u7eaa\u8981", 1.4, MUTED)
    add_text(slide, 150, 92.5, 20, 4, page, 1.4, MUTED, align=PP_ALIGN.RIGHT, font=TITLE_FONT)


def cover_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)

    # Upper-left brand marker
    add_rect(slide, 8, 8, 2.0, 2.0, fill=RED)
    add_text(slide, 11, 7.5, 60, 5, "\u793a\u4f8b\u8bfe\u7a0b \u00b7 \u7b2c\u4e00\u8bb2", 1.9, MUTED, font=TITLE_FONT, bold=True)

    # Upper-right decorative seal; no external SVG/PNG dependency

    # Editorial cover titles
    add_text(slide, 8, 20, 85, 10, "\u793a\u4f8b\u4e3b\u9898\u7684", 7.5, INK, bold=True, font=TITLE_FONT)
    add_text(slide, 8, 31, 85, 10, "\u6458\u8981\u6807\u9898", 7.5, RED, bold=True, font=TITLE_FONT)
    
    # Red accent highlight bar
    add_rect(slide, 8, 42.5, 30, 0.6, fill=RED)
    
    add_text(slide, 8, 44, 85, 10, "\u4e0e\u526f\u9898", 7.5, INK, bold=True, font=TITLE_FONT)
    
    # Ribbon stamp
    add_rect(slide, 8, 57, 40, 7.8, fill=RED)
    add_text(slide, 8, 58.5, 40, 5, "\u8bfe\u7a0b\u7eaa\u8981", 2.6, ON_ACCENT, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    
    # Elegant metadata block in card layout with high contrast ink text
    add_rect(slide, 8, 68, 80, 10.5, fill=CARD, line=DIVIDER, line_w=0.75)
    add_text(slide, 10, 69.2, 76, 8, "\u65e5\u671f  2026-01-01   |   \u65f6\u957f  \u7ea6 1 \u5c0f\u65f6\n\u4e3b\u8bb2  \u4e3b\u8bb2\u4eba\u7532   |   \u53c2\u4e0e  \u53c2\u4e0e\u4eba\u4e59 \u00b7 \u53c2\u4e0e\u4eba\u4e19", 1.9, INK, font=BODY_FONT, bold=True)

    # Core concept quote card
    add_rect(slide, 94, 63, 76, 17, fill=CARD)
    add_rect(slide, 94, 63, 0.8, 17, fill=RED)
    add_text(slide, 97, 65, 70, 4, "CORE THESIS  \u4e00\u53e5\u8bdd\u4e3b\u65e8", 1.5, RED, bold=True, font=TITLE_FONT)
    add_text(slide, 97, 70, 70, 9, "\u628a\u201c\u6838\u5fc3\u89c2\u70b9\u201d\u7528\u4e00\u53e5\u8bdd\u8bb2\u6e05\u695a\uff0c\u8fd9\u91cc\u662f\u793a\u4f8b\u91d1\u53e5\u3002", 2.2, INK, bold=True, font=BODY_FONT)
    
    footer(slide, next_page())

def info_slides():
    rows = [
        ("\u8bfe\u7a0b", "\u793a\u4f8b\u8bfe\u7a0b \u00b7 \u7b2c\u4e00\u8bb2", None),
        ("\u4e3b\u9898", "\u793a\u4f8b\u4e3b\u9898\u7684\u6458\u8981\u6807\u9898\u4e0e\u526f\u9898", None),
        ("\u4e3b\u8bb2", "\u4e3b\u8bb2\u4eba\u7532", "\u804c\u4e1a\u793a\u4f8b"),
        ("\u53c2\u4e0e", "\u53c2\u4e0e\u4eba\u4e59 \u00b7 \u53c2\u4e0e\u4eba\u4e19", "\u53c2\u4e0e\u4eba\u4e59\u8bb2\u6a21\u5757\u4e00 / \u53c2\u4e0e\u4eba\u4e19\u8bb2\u6a21\u5757\u4e8c"),
        ("\u5f62\u5f0f", "\u4e3b\u8bb2 + \u5b66\u5458\u95ee\u7b54 + \u5de5\u5177\u5b9e\u64cd", None),
        ("\u914d\u5957\u5de5\u5177", "\u793a\u4f8b\u5de5\u5177", "\u542b\u4f53\u9a8c\u989d\u5ea6"),
    ]
    for part in chunk(rows, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, PAPER)
        header(slide, "\u58f9", "\u8bfe\u7a0b\u4fe1\u606f", "OVERVIEW")
        
        # Border top
        add_rect(slide, 5, 24, 168, 0.13, fill=DIVIDER)
        y = 26
        for label, value, sub in part:
            add_text(slide, 5, y, 20, 5, label, 1.75, MUTED, font=TITLE_FONT, bold=True)
            add_text(slide, 29, y, 138, 6, value, 2.4, INK, bold=True)
            if sub:
                add_text(slide, 29, y + 4.2, 138, 4, sub, 1.9, MUTED)
            add_rect(slide, 5, y + 9.5, 168, 0.13, fill=DIVIDER)
            y += 10.5
        footer(slide, next_page())
def timeline_slides():
    steps = [
        ("\u4ece\u4e00\u4e2a\u793a\u4f8b\u573a\u666f\u5207\u5165", "\u793a\u4f8b\u573a\u666f\u63cf\u8ff0\uff1a\u8fd9\u91cc\u8bf4\u660e\u7b2c\u4e00\u4e2a\u8981\u70b9\uff0c\u7528\u4e8e\u6f14\u793a\u65f6\u95f4\u7ebf\u7ec4\u4ef6\u3002", "1"),
        ("\u629b\u51fa\u6838\u5fc3\u89c2\u70b9", "\u793a\u4f8b\u6982\u5ff5\uff1a\u8fd9\u91cc\u5c55\u793a\u4e00\u4e2a\u5173\u952e\u6982\u5ff5\uff1b\u53e6\u4e00\u4e2a\u6982\u5ff5\uff1a\u8fd9\u91cc\u5c55\u793a\u4e0e\u4e4b\u547c\u5e94\u7684\u8bf4\u660e\u3002", "2"),
        ("\u7ed9\u51fa\u7ed3\u8bba", "\u505aA\u7c7b\u4e8b\u60c5\u7684\u5c06\u5982\u4f55\uff1b\u6709B\u7c7b\u89c6\u89d2\u7684\u4eba\u5c06\u5982\u4f55\u3002", "3"),
        ("\u793a\u4f8b\u8981\u70b9", "\u793a\u4f8b\u8bf4\u660e\uff1a\u8fd9\u91cc\u5c55\u793a\u65f6\u95f4\u7ebf\u7ec4\u4ef6\u7684\u7b2c\u56db\u6b65\u6392\u5e03\u6548\u679c\u3002", "4"),
    ]
    for part in chunk(steps, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, PAPER)
        header(slide, "\u8d30", "\u601d\u8def\u5f15\u5bfc", "NARRATIVE FLOW")
        add_rect(slide, 8, 27, 0.13, 58, fill=DIVIDER)
        y = 27
        for title, body, num in part:
            add_rect(slide, 5, y, 4.2, 4.2, fill=PAPER, line=DIVIDER, line_w=0.75)
            add_text(slide, 5, y, 4.2, 4.2, num, 1.9, RED, align=PP_ALIGN.CENTER, font=TITLE_FONT, bold=True)
            add_text(slide, 13, y, 150, 5, title, 2.7, INK, bold=True, font=TITLE_FONT)
            add_text(slide, 13, y + 5.5, 150, 8, body, 2.1, NARRATIVE)
            y += 19.5
        footer(slide, next_page())


def viewpoint_slides():
    quotes = [
        "\u89c2\u70b9\u4e00\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e00\u6761\u793a\u4f8b\u89c2\u70b9\uff0c\u7528\u4e8e\u6f14\u793a\u89c2\u70b9\u5361\u5e03\u5c40\u3002",
        "\u89c2\u70b9\u4e8c\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e8c\u6761\u793a\u4f8b\u89c2\u70b9\uff0c\u8bf4\u660e\u89c2\u70b9\u5361\u7684\u6392\u5e03\u65b9\u5f0f\u3002",
        "\u89c2\u70b9\u4e09\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e09\u6761\u793a\u4f8b\u89c2\u70b9\uff0c\u6f14\u793a\u5206\u7ec4\u95f4\u8ddd\u6548\u679c\u3002",
        "\u5174\u8da3\u9a71\u52a8\u578b\u5b66\u4e60\u662f \u793a\u4f8b\u4e3b\u9898\u7684\u5b66\u4e60\u65b9\u5f0f\uff0c\u5185\u9a71\u529b\u8fdc\u80dc\u5916\u9a71\u529b\u3002",
    ]
    for part in chunk(quotes, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, PAPER)
        header(slide, "\u53c1", "\u8bb2\u8005\u89c2\u70b9", "VIEWPOINTS")
        
        # Red vertical indicator line for clean sidebar layout
        add_rect(slide, 5, 25, 0.6, 12, fill=RED)
        add_text(slide, 8, 25, 45, 4, "\u4e3b\u8bb2 \u00b7 \u793a\u4f8b\u89d2\u8272\u5b9a\u4f4d", 1.5, RED, font=TITLE_FONT, bold=True)
        add_text(slide, 8, 29, 45, 6, "\u4e3b\u8bb2\u4eba\u7532", 2.6, INK, bold=True, font=TITLE_FONT)
        
        y = 26
        for text in part:
            # Elegant card containers for viewpoints, complete with borders and left red highlights
            add_rect(slide, 58, y, 114, 11.5, fill=CARD, line=DIVIDER, line_w=0.75)
            add_rect(slide, 58, y, 0.8, 11.5, fill=RED)
            add_text(slide, 61, y + 1.2, 5, 5, "\u201c", 4.0, RED, font=TITLE_FONT, bold=True)
            add_text(slide, 66, y + 1.8, 102, 8, text, 2.3, BODY)
            y += 13.8
        footer(slide, next_page())
def action_slides():
    items = [
        "\u884c\u52a8\u9879\u4e00\uff1a\u8fd9\u91cc\u5199\u7b2c\u4e00\u6761\u884c\u52a8\u5efa\u8bae\uff0c\u7528\u4e8e\u6f14\u793a\u6e05\u5355\u7ec4\u4ef6\u3002",
        "\u884c\u52a8\u9879\u4e8c\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e8c\u6761\u884c\u52a8\u5efa\u8bae\u7684\u8bf4\u660e\u6587\u5b57\u3002",
        "\u884c\u52a8\u9879\u4e09\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e09\u6761\u884c\u52a8\u5efa\u8bae\uff0c\u5c55\u793a\u6e05\u5355\u7684\u6392\u5e03\u3002",
        "\u884c\u52a8\u9879\u56db\uff1a\u8fd9\u91cc\u662f\u7b2c\u56db\u6761\u884c\u52a8\u5efa\u8bae\uff0c\u6f14\u793a\u6e05\u5355\u7684\u6536\u5c3e\u3002",
        "\u884c\u52a8\u9879\u4e94\uff1a\u8fd9\u91cc\u662f\u7b2c\u4e94\u6761\u884c\u52a8\u5efa\u8bae\uff0c\u6f14\u793a\u6e05\u5355\u7684\u6269\u5c55\u3002",
    ]
    parts = chunk(items, 4)
    for idx, part in enumerate(parts):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, PAPER)
        header(slide, "\u8086", "\u7ed9\u5b66\u5458\u7684\u884c\u52a8\u6e05\u5355", "ACTION ITEMS")
        
        # White container table card
        add_rect(slide, 5, 24, 168, 48, fill=RGBColor(255, 255, 255), line=DIVIDER, line_w=1.0)
        y = 26
        for text in part:
            add_rect(slide, 8, y, 3.0, 3.0, fill=None, line=RED, line_w=1.5)
            add_rect(slide, 9.2, y + 0.8, 0.6, 1.4, fill=RED).rotation = 45
            add_text(slide, 13, y - 0.2, 155, 6, text, 2.3, INK)
            if text != part[-1]:
                add_rect(slide, 8, y + 8.2, 162, 0.13, fill=DIVIDER)
            y += 10.5
            
        if idx == len(parts) - 1:
            # note-box card container with left red strip
            add_rect(slide, 5, 75, 168, 11, fill=RGBColor(255, 255, 255), line=DIVIDER, line_w=0.75)
            add_rect(slide, 5, 75, 0.8, 11, fill=RED)
            add_text(slide, 8, 77.2, 160, 7, "\u8bf4\u660e\uff1a\u672c\u6a21\u677f\u6309\u6696\u7eb8\u793a\u4f8b\u81ea\u9002\u5e94\u89c4\u8303\u751f\u6210\uff0c\u6b63\u6587\u4e3a\u53ef\u7f16\u8f91\u6587\u672c\u6846\u3002", 1.9, MUTED)
        footer(slide, next_page())
def colophon_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PAPER)
    
    # Outer elegant thin border for certificate colophon style
    add_rect(slide, 5, 5, 168, 81, fill=None, line=DIVIDER, line_w=1.5)
    
    # Center Red Seal
    add_rect(slide, 83, 18, 14, 14, fill=RED, shape=MSO_SHAPE.OVAL)
    add_text(slide, 83, 22.5, 14, 6, "\u793a\u4f8b", 3.2, ON_ACCENT, bold=True, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    
    # Brand details
    add_text(slide, 30, 38, 120, 6, "\u793a\u4f8b\u8bfe\u7a0b \u00b7 \u7b2c\u4e00\u8bb2", 3.0, INK, bold=True, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_text(slide, 30, 46, 120, 5, "\u8bfe\u7a0b\u7eaa\u8981   \u00b7   \u793a\u4f8b\u5185\u5bb9", 1.8, MUTED, align=PP_ALIGN.CENTER)
    
    # Separation line
    add_rect(slide, 40, 54, 100, 0.15, fill=RED)
    
    # Description block
    add_rect(slide, 25, 60, 130, 15, fill=CARD, line=DIVIDER, line_w=0.75)
    add_text(slide, 28, 62, 124, 11, "\u8bf4\u660e\uff1a\u672c\u9875\u4e3a\u6a21\u677f\u793a\u4f8b\u5185\u5bb9\uff0c\u66ff\u6362\u4e3a\u4f60\u81ea\u5df1\u7684\u6b63\u6587\u5373\u53ef\u3002\u6b63\u6587\u4e3a\u53ef\u7f16\u8f91\u6587\u672c\u6846\u3002", 2.0, NARRATIVE, font=BODY_FONT)
    
    footer(slide, next_page())

def main():
    out = sys.argv[1] if len(sys.argv) > 1 else "\u793a\u4f8b-\u8f93\u51fa.pptx"
    Path(out).parent.mkdir(parents=True, exist_ok=True)
    cover_slide()
    info_slides()
    timeline_slides()
    viewpoint_slides()
    action_slides()
    colophon_slide()
    prs.save(out)
    print("saved", out)
    print("slides", slide_no)


if __name__ == "__main__":
    run_main(main, next_step="请根据上面的提示处理；缺少依赖时运行 install.bat（Windows）或 install.sh（macOS/Linux）")
